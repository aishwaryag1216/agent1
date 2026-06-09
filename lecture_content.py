import streamlit as st
from google import genai
from docx import Document
from pptx import Presentation
import io

# =========================
# Configure Gemini API
# =========================

API_KEY = st.secrets["GEMINI_API_KEY1"]
client = genai.Client(api_key=API_KEY)

# =========================
# Streamlit UI
# =========================
st.set_page_config(page_title="Lecture Content Generator")

st.title("📚 Lecture Content Generator")

# Input Fields
subject = st.text_input("Enter Subject Name")
topic = st.text_input("Enter Topic Name")

# Output Format Dropdown
output_format = st.selectbox(
    "Select Output Format",
    ["Text", "Document", "PPT"]
)

# =========================
# Generate Content Function
# =========================
def generate_content(subject, topic):
    try:
        prompt = f"""
        Subject: {subject}
        Topic: {topic}

        Generate detailed lecture content.
        """

        response = client.models.generate_content(
            model="gemini-2.5-pro",
            contents=prompt
        )

        return response.text

    except Exception as e:
        st.error(f"Gemini Error: {e}")
        return None

  

        return response.text 
       

# =========================
# Create DOCX
# =========================
def create_docx(subject, topic, content):
    doc = Document()

    doc.add_heading(subject, level=1)
    doc.add_heading(topic, level=2)
    doc.add_paragraph(content)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    return buffer

# =========================
# Create PPT
# =========================
def create_ppt(subject, topic, content):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = topic
    slide.placeholders[1].text = content[:1000]

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer

# =========================
# Generate Button
# =========================
if st.button("Generate"):

    if not subject or not topic:
        st.warning("Please enter Subject Name and Topic Name.")
    else:
        with st.spinner("Generating content..."):
            content = generate_content(subject, topic)

        # Text Output
        if output_format == "Text":
            st.subheader("Generated Content")
            st.write(content)

        # DOCX Output
        elif output_format == "Document":
            doc_file = create_docx(subject, topic, content)

            st.download_button(
                label="📄 Download Document",
                data=doc_file,
                file_name=f"{topic}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )

        # PPT Output
        elif output_format == "PPT":
            ppt_file = create_ppt(subject, topic, content)

            st.download_button(
                label="📊 Download PPT",
                data=ppt_file,
                file_name=f"{topic}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
